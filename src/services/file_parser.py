"""多格式文件解析器"""
from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from src.models.parsed_content import StructuredBlock


def _read_file_text(path: Path) -> str:
    """自动检测编码并读取文本文件"""
    raw = path.read_bytes()
    if not raw:
        return ""
    from charset_normalizer import from_bytes
    result = from_bytes(raw)
    if result and result.best():
        return str(result.best())
    return raw.decode("utf-8", errors="replace")


@dataclass
class ParsedFile:
    title: str
    content: str
    file_type: str
    source_path: str
    metadata: dict
    structured: list["StructuredBlock"] | None = None


def parse_file(file_path: str) -> list[ParsedFile]:
    """解析文件，返回 ParsedFile 列表（大多数格式返回单元素列表，Excel 返回多元素）"""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    ext = path.suffix.lower()
    parsers = {
        ".pdf": _parse_pdf,
        ".pptx": _parse_pptx,
        ".ppt": _parse_pptx,
        ".docx": _parse_docx,
        ".doc": _parse_docx,
        ".xlsx": _parse_excel,
        ".xls": _parse_excel,
        ".csv": _parse_csv,
        ".txt": _parse_text,
        ".md": _parse_markdown,
        ".html": _parse_html,
        ".htm": _parse_html,
    }
    code_extensions = {
        ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".go", ".rs",
        ".rb", ".php", ".cs", ".swift", ".kt", ".scala", ".sh", ".bat",
        ".sql", ".r", ".m", ".json", ".yaml", ".yml", ".xml", ".toml",
    }

    if ext in parsers:
        result = parsers[ext](path)
        # Excel 解析器直接返回 list[ParsedFile]
        if isinstance(result, list):
            return result
        return [result]
    elif ext in code_extensions:
        return [_parse_code(path)]
    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp"):
        return [_parse_image(path)]
    else:
        return [_parse_text(path)]


MAX_ROWS_PER_SHEET = 500


def _looks_like_data(value) -> bool:
    s = str(value).strip()
    if not s:
        return False
    try:
        float(s)
        return True
    except ValueError:
        pass
    import re
    if re.match(r'^\d{4}[-/]\d{1,2}[-/]\d{1,2}$', s):
        return True
    return False


def _detect_headers(first_row: list) -> list[str]:
    non_empty = [v for v in first_row if str(v).strip()]
    if not non_empty:
        return []
    all_stringish = all(
        not _looks_like_data(v)
        for v in non_empty
    )
    if all_stringish:
        return [str(v).strip() for v in first_row]
    return []


def _table_rows_to_text(headers: list[str], rows: list[list[str]], sheet_name: str = "") -> str:
    lines = []
    if sheet_name:
        lines.append(f"[工作表: {sheet_name}]")
    if headers:
        lines.append(f"表头: {' | '.join(headers)}")
    lines.append("")

    total = len(rows)
    truncated = total > MAX_ROWS_PER_SHEET
    row_limit = min(total, MAX_ROWS_PER_SHEET)

    for i in range(row_limit):
        row = rows[i]
        parts = []
        for col_idx, val in enumerate(row):
            val = str(val).strip() if val is not None else ""
            if not val:
                continue
            header = headers[col_idx] if col_idx < len(headers) else f"列{col_idx + 1}"
            parts.append(f"{header}={val}")
        if not parts:
            continue
        lines.append(f"第{i + 1}行: {', '.join(parts)}")

    if truncated:
        lines.append(f"\n（共{total}行，已截取前{MAX_ROWS_PER_SHEET}行）")

    return "\n".join(lines)


def _table_rows_to_blocks(
    headers: list[str],
    rows: list[list[str]],
    sheet_name: str = "",
) -> list["StructuredBlock"]:
    """将表格数据转换为结构化 Block 树

    每行数据变成一个命名 Block（用首列非空值命名），
    各列值作为子 Block（property 格式: 列名:: 值）。
    表头信息保存在每个行 Block 的 properties 中，确保检索时上下文不丢失。
    """
    from src.models.parsed_content import StructuredBlock

    total = len(rows)
    row_limit = min(total, MAX_ROWS_PER_SHEET)
    col_names = headers if headers else []
    truncated = total > MAX_ROWS_PER_SHEET

    blocks: list[StructuredBlock] = []

    for i in range(row_limit):
        row = rows[i]
        # 跳过完全空行
        non_empty = [(col_idx, str(v).strip()) for col_idx, v in enumerate(row) if str(v).strip()]
        if not non_empty:
            continue

        # 首列非空值作为 Block 名称（便于识别）
        first_val = non_empty[0][1] if non_empty else f"第{i + 1}行"
        row_content = f"**{first_val}**"

        # 行 Block 的属性：携带表头上下文
        row_props: dict = {"row_index": i + 1}
        if col_names:
            row_props["columns"] = " | ".join(col_names)
        if sheet_name:
            row_props["sheet"] = sheet_name
        row_parts = []
        for col_idx, val in non_empty:
            header = col_names[col_idx] if col_idx < len(col_names) else f"col{col_idx + 1}"
            row_parts.append(f"{header}={val}")
        if row_parts:
            row_content = f"{non_empty[0][1]}: " + " | ".join(row_parts)

        row_block = StructuredBlock(
            content=row_content,
            block_type="table_row",
            properties=row_props,
        )

        # 每个非空列值作为子 Block（property 格式）
        for col_idx, val in non_empty:
            header = col_names[col_idx] if col_idx < len(col_names) else f"列{col_idx + 1}"
            row_block.children.append(StructuredBlock(
                content=f"{header}:: {val}",
                block_type="property",
                properties={"column": header, "value": val},
            ))

        blocks.append(row_block)

    if truncated:
        # 截断提示作为最后一个 Block
        blocks.append(StructuredBlock(
            content=f"（共{total}行，已截取前{MAX_ROWS_PER_SHEET}行）",
            block_type="text",
            properties={"truncated": True, "total_rows": total},
        ))

    return blocks


def _extract_pptx_shapes(shapes) -> list[str]:
    """递归提取形状中的文本（含组合形状）"""
    parts = []
    for shape in shapes:
        # 组合形状：递归提取子形状
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                parts.extend(_extract_pptx_shapes(shape.shapes))
            except Exception:
                logger.debug("Failed to extract group shape, skipping")
            continue
        # 普通文本框/标题
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
        # 表格
        if shape.has_table:
            rows = []
            for row in shape.table.rows:
                cells = [cell.text_frame.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            table_text = "\n".join(rows)
            if table_text.strip():
                parts.append(table_text)
    return parts


def _pptx_slide_to_block(slide_num: int, parts: list[str]) -> "StructuredBlock":
    """将 PPT 幻灯片内容转换为结构化 Block

    每张幻灯片 → 一个 slide Block，各文本部分 → 子 Block。
    第一个短文本视为幻灯片标题，其余视为正文。
    """
    from src.models.parsed_content import StructuredBlock

    # 第一个短文本作为标题
    title_text = ""
    body_parts = []
    for p in parts:
        if not title_text and len(p.strip()) <= 80:
            title_text = p.strip()
        else:
            body_parts.append(p)

    slide_block = StructuredBlock(
        content=f"[第{slide_num}页] {title_text}" if title_text else f"[第{slide_num}页]",
        block_type="slide",
        properties={"slide": slide_num},
        level=0,
    )

    for part in body_parts:
        slide_block.children.append(StructuredBlock(
            content=part,
            block_type="text",
        ))

    return slide_block


def _parse_pptx(path: Path) -> ParsedFile:
    import re as _re

    from pptx import Presentation

    prs = Presentation(str(path))
    slide_count = len(prs.slides)
    slides = []
    structured_blocks: list = []

    for i, slide in enumerate(prs.slides):
        parts = _extract_pptx_shapes(slide.shapes)

        # 备注页
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[备注] {notes}")

        if parts:
            cleaned = []
            for p in parts:
                # 清理控制字符和特殊项目符号
                p = _re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", p)
                p = p.replace("\r\n", "\n").strip()
                if p:
                    cleaned.append(p)
            if cleaned:
                slides.append(f"[第{i+1}页]\n" + "\n\n".join(cleaned))

                # 结构化：每张幻灯片 → 一个 Block
                slide_block = _pptx_slide_to_block(i + 1, cleaned)
                structured_blocks.append(slide_block)

    content = "\n\n".join(slides)

    # 从首页提取标题（首个非空短文本）
    title = path.stem
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and len(t) <= 80:
                    title = t
                    break

    return ParsedFile(
        title=title,
        content=content,
        file_type="pptx",
        source_path=str(path),
        metadata={"slides": slide_count},
        structured=structured_blocks if structured_blocks else None,
    )


def _pdf_pages_to_blocks(pdf_pages) -> list["StructuredBlock"]:
    """将 PDF 页面转换为结构化 Block 树

    每页 → heading Block，段落 → 子 Block。
    标题启发式检测：短行且不以句号结尾的视为 heading。
    """
    from src.models.parsed_content import StructuredBlock

    blocks: list[StructuredBlock] = []

    for i, page in enumerate(pdf_pages):
        text = page.extract_text()
        if not text or not text.strip():
            continue

        page_block = StructuredBlock(
            content=f"[第{i + 1}页]",
            block_type="heading",
            properties={"page": i + 1},
            level=0,
        )

        # 按段落分割
        for para in text.split("\n\n"):
            para = para.strip()
            if not para:
                continue
            # 按行进一步细分长段落
            lines = para.split("\n")
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                # 启发式标题检测：短行 + 不以标点结尾
                is_heading = (
                    len(line) < 60
                    and not line.endswith(("。", "，", "；", "：", "！", "？", ".", ",", ";"))
                    and not line.startswith(("—", "－", "-"))
                )
                page_block.children.append(StructuredBlock(
                    content=line,
                    block_type="heading" if is_heading else "text",
                    level=1 if is_heading else 2,
                ))

        if page_block.children:
            blocks.append(page_block)

    return blocks


_PDF_WATERMARK_TOKEN_RE = None


def _pdf_watermark_token_re():
    global _PDF_WATERMARK_TOKEN_RE
    if _PDF_WATERMARK_TOKEN_RE is None:
        import re
        keywords = [
            "confidential", "draft", "internal", "sample",
            "unauthorized", "proprietary",
            "\u6c34\u5370", "\u5185\u90e8\u6587\u4ef6", "\u4ec5\u4f9b\u53c2\u8003",
            "\u673a\u5bc6", "\u6837\u672c", "\u4e25\u7981\u590d\u5236",
            "\u4e0d\u5f97\u5916\u4f20", "\u672a\u7ecf\u6388\u6743",
        ]
        _PDF_WATERMARK_TOKEN_RE = re.compile(
            r"(?<![\w\u4e00-\u9fff])(?:"
            + "|".join(re.escape(k) for k in keywords)
            + r")(?![\w\u4e00-\u9fff])",
            re.IGNORECASE,
        )
    return _PDF_WATERMARK_TOKEN_RE


def _pdf_line_norm(line: str) -> str:
    return "".join(line.strip().lower().split())


def _pdf_repeated_watermark_norms(pages_text: list[str]) -> set[str]:
    if len(pages_text) < 3:
        return set()
    counts: dict[str, int] = {}
    for text in pages_text:
        seen = {_pdf_line_norm(line) for line in text.split("\n") if line.strip()}
        for norm in seen:
            if norm:
                counts[norm] = counts.get(norm, 0) + 1
    threshold = len(pages_text) * 0.8
    return {norm for norm, count in counts.items() if count >= threshold}


def _clean_pdf_watermark_line(line: str, repeated_norms: set[str] | None = None) -> str:
    if not line.strip():
        return ""
    if repeated_norms and _pdf_line_norm(line) in repeated_norms:
        return ""
    cleaned = _pdf_watermark_token_re().sub("", line)
    return " ".join(cleaned.split())


def _remove_pdf_watermarks(pages_text: list[str]) -> list[str]:
    """去除 PDF 水印文本

    策略：出现在 >80% 页面的相同文本行，极大概率是水印或页眉页脚。
    """
    repeated_norms = _pdf_repeated_watermark_norms(pages_text)
    cleaned_pages = []
    for text in pages_text:
        lines = []
        for line in text.split("\n"):
            cleaned = _clean_pdf_watermark_line(line, repeated_norms)
            if cleaned.strip():
                lines.append(cleaned)
        cleaned_pages.append("\n".join(lines))
    return cleaned_pages


def _filter_watermark_keywords(text: str) -> str:
    """过滤文本中的常见水印关键词行（用于页数太少无法统计的 PDF）"""
    import re
    _WATERMARK_RE = re.compile(
        r'^(?:水印|confidential|draft|internal|内部文件|仅供参考|机密|样本|sample|'
        r'严禁复制|不得外传|未经授权|unauthorized|proprietary)$',
        re.IGNORECASE,
    )
    lines = text.split("\n")
    filtered = [line for line in lines if not _WATERMARK_RE.match(line.strip())]
    return "\n".join(filtered)


def _pdf_pages_to_blocks_from_text(pages_text: list[str]) -> list["StructuredBlock"]:
    """从纯文本列表构建结构化 Block 树（替代直接依赖 reader.pages）

    用于已经过水印清洗的文本，避免二次 extract_text() 调用。
    """
    from src.models.parsed_content import StructuredBlock

    blocks: list[StructuredBlock] = []

    for i, text in enumerate(pages_text):
        if not text or not text.strip():
            continue

        page_block = StructuredBlock(
            content=f"[第{i + 1}页]",
            block_type="heading",
            properties={"page": i + 1},
            level=0,
        )

        for para in text.split("\n\n"):
            para = para.strip()
            if not para:
                continue
            lines = para.split("\n")
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                is_heading = (
                    len(line) < 60
                    and not line.endswith(("。", "，", "；", "：", "！", "？", ".", ",", ";"))
                    and not line.startswith(("—", "－", "-"))
                )
                page_block.children.append(StructuredBlock(
                    content=line,
                    block_type="heading" if is_heading else "text",
                    level=1 if is_heading else 2,
                ))

        if page_block.children:
            blocks.append(page_block)

    return blocks


def _parse_pdf(path: Path) -> ParsedFile:
    from io import BytesIO

    from PyPDF2 import PdfReader
    from PyPDF2.errors import PdfReadError, PdfStreamError

    # ---- 尝试打开 PDF，失败时用 pikepdf 修复 ----
    reader = None
    try:
        reader = PdfReader(str(path))
    except (PdfReadError, PdfStreamError, Exception) as e:
        err_msg = str(e).lower()
        if "eof" in err_msg or "stream" in err_msg or "marker" in err_msg:
            # 尝试用 pikepdf 修复损坏的 PDF
            try:
                import pikepdf
                with pikepdf.open(str(path), repair=True) as pdf:
                    buf = BytesIO()
                    pdf.save(buf)
                    buf.seek(0)
                reader = PdfReader(buf)
            except ImportError:
                raise ValueError(
                    f"PDF 文件可能已损坏 ({path.name}): {e}。"
                    f"安装 pikepdf 可尝试自动修复: pip install pikepdf"
                )
            except Exception as repair_err:
                raise ValueError(
                    f"PDF 文件已损坏且无法修复 ({path.name}): {repair_err}"
                )
        else:
            raise ValueError(f"无法读取 PDF 文件 ({path.name}): {e}")

    # ---- 加密 PDF 处理 ----
    if reader.is_encrypted:
        decrypt_result = reader.decrypt("")
        if decrypt_result == 0:
            try:
                import pikepdf
                with pikepdf.open(str(path)) as pdf:
                    buf = BytesIO()
                    pdf.save(buf)
                    buf.seek(0)
                reader = PdfReader(buf)
            except ImportError:
                raise ValueError(
                    f"PDF 文件已加密 ({path.name})，"
                    f"需要安装 pikepdf 才能导入: pip install pikepdf"
                )
            except Exception:
                raise ValueError(
                    f"PDF 文件已加密且需要密码才能打开 ({path.name})，"
                    f"请先解密或提供密码后重试"
                )

    # ---- 提取文本 + 去水印 ----
    raw_pages: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text()
        except Exception:
            # 单页提取失败不影响其他页
            logger.debug("Failed to extract text from PDF page %d", i)
            text = ""
        if text:
            raw_pages.append(text)

    # 去除水印文本
    cleaned_pages = _remove_pdf_watermarks(raw_pages)

    pages = []
    for i, text in enumerate(cleaned_pages):
        if text.strip():
            pages.append(f"[第{i+1}页]\n{text}")
    content = "\n\n".join(pages)

    metadata = {"pages": len(reader.pages)}
    if reader.is_encrypted:
        metadata["encrypted"] = True

    # 结构化 Block 树：使用去水印后的文本
    structured = _pdf_pages_to_blocks_from_text(cleaned_pages)

    return ParsedFile(
        title=path.stem,
        content=content,
        file_type="pdf",
        source_path=str(path),
        metadata=metadata,
        structured=structured,
    )


def _docx_to_blocks(doc) -> list["StructuredBlock"]:
    """将 Word 文档转换为结构化 Block 树

    利用 python-docx 的段落样式（Heading 1/2/3）建立层级。
    表格复用 _table_rows_to_blocks 结构化逻辑。
    """
    from src.models.parsed_content import StructuredBlock

    blocks: list[StructuredBlock] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""
        if "Heading" in style_name or "标题" in style_name:
            # 从样式名中提取标题级别
            level = 1
            for ch in reversed(style_name):
                if ch.isdigit():
                    level = int(ch)
                    break
            blocks.append(StructuredBlock(
                content=text,
                block_type="heading",
                level=level,
                properties={"style": style_name},
            ))
        else:
            blocks.append(StructuredBlock(
                content=text,
                block_type="text",
            ))

    # Word 表格：复用表格结构化逻辑
    for table in doc.tables:
        rows_data = []
        for row in table.rows:
            rows_data.append([cell.text.strip() for cell in row.cells])
        if not rows_data:
            continue

        headers = _detect_headers(rows_data[0])
        data_rows = rows_data[1:] if headers else rows_data

        # 表格作为一个整体 Block
        table_block = StructuredBlock(
            content="[表格]",
            block_type="text",
            properties={"type": "table", "columns": " | ".join(headers) if headers else ""},
        )
        table_block.children = _table_rows_to_blocks(headers, data_rows)
        blocks.append(table_block)

    return blocks


def _parse_docx(path: Path) -> ParsedFile:
    from docx import Document
    doc = Document(str(path))

    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            paragraphs.append(" | ".join(cells))
    content = "\n\n".join(paragraphs)

    # 结构化 Block 树：利用段落样式建立层级，表格结构化
    structured = _docx_to_blocks(doc)

    return ParsedFile(
        title=path.stem,
        content=content,
        file_type="docx",
        source_path=str(path),
        metadata={},
        structured=structured,
    )


def _parse_text(path: Path) -> ParsedFile:
    content = _read_file_text(path)
    return ParsedFile(
        title=path.stem,
        content=content,
        file_type="txt",
        source_path=str(path),
        metadata={},
    )


def _parse_markdown(path: Path) -> ParsedFile:
    content = _read_file_text(path)
    return ParsedFile(
        title=path.stem,
        content=content,
        file_type="md",
        source_path=str(path),
        metadata={},
    )


def _parse_html(path: Path) -> ParsedFile:
    from bs4 import BeautifulSoup
    content = _read_file_text(path)
    soup = BeautifulSoup(content, "html.parser")
    title = soup.title.string.strip() if soup.title and soup.title.string else path.stem
    text = _extract_main_content(soup)
    return ParsedFile(
        title=title,
        content=text,
        file_type="html",
        source_path=str(path),
        metadata={},
    )


def _extract_main_content(soup) -> str:
    """从 HTML 中提取正文区域，优先定位 <article>/<main>，回退全文"""
    REMOVE_TAGS = {"script", "style", "nav", "footer", "header", "aside",
                   "form", "iframe", "noscript", "svg", "button", "input",
                   "select", "textarea"}
    REMOVE_ROLES = {"navigation", "banner", "complementary", "contentinfo",
                    "search", "toolbar", "menu", "menubar"}

    # 先移除无关标签
    for tag in soup.find_all(lambda t: t.name in REMOVE_TAGS):
        tag.decompose()
    for tag in soup.find_all(lambda t: t.get("role", "") in REMOVE_ROLES):
        tag.decompose()
    # 移除广告、评论区等常见非正文容器
    for cls_pattern in ["ad", "ads", "advert", "comment", "sidebar",
                        "social", "share", "related", "recommend", "widget",
                        "cookie", "popup", "modal", "banner", "toolbar"]:
        for tag in soup.find_all(class_=lambda c: c and cls_pattern in " ".join(c).lower()):
            tag.decompose()

    # 优先定位正文容器
    main = None
    for selector in [
        lambda s: s.find("article"),
        lambda s: s.find("main"),
        lambda s: s.find(attrs={"role": "main"}),
        lambda s: s.find(class_=lambda c: c and any(
            k in " ".join(c).lower()
            for k in ["article", "post", "content", "entry", "body"]
        )),
        lambda s: s.find(id=lambda i: i and any(
            k in i.lower()
            for k in ["article", "post", "content", "entry", "body"]
        )),
    ]:
        candidate = selector(soup)
        if candidate:
            text = candidate.get_text(separator="\n", strip=True)
            if len(text) > 200:
                main = candidate
                break

    if main:
        text = main.get_text(separator="\n", strip=True)
    else:
        text = soup.get_text(separator="\n", strip=True)

    # 清理：合并连续空行
    import re
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def parse_url(url: str, timeout: float | None = None) -> ParsedFile:
    """抓取网页并提取正文文本"""
    import ipaddress
    import re
    import socket
    from urllib.parse import urlparse

    import httpx
    from bs4 import BeautifulSoup

    if not url.startswith(("http://", "https://")):
        raise ValueError(f"不支持的 URL 协议: {url}")

    # SSRF 防护：阻止对内网/回环/链路本地地址的请求
    parsed = urlparse(url)
    hostname = parsed.hostname
    if hostname:
        try:
            # 先尝试 DNS 解析，再检查 IP 是否为私有地址
            resolved_ips = socket.getaddrinfo(hostname, None, socket.AF_UNSPEC, socket.SOCK_STREAM)
            for _, _, _, _, addr in resolved_ips:
                ip = ipaddress.ip_address(addr[0])
                if ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_reserved:
                    raise ValueError(f"不允许访问内网地址: {hostname} ({ip})")
        except socket.gaierror:
            raise ValueError(f"无法解析主机名: {hostname}")

    if timeout is None:
        from src.utils.config import Config
        timeout = Config.get("web.timeout", 30)

    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                      "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
    }
    with httpx.Client(timeout=timeout, follow_redirects=True, headers=headers, max_redirects=5) as client:
        response = client.get(url)
        if response.status_code != 200:
            raise RuntimeError(f"HTTP {response.status_code}: {url}")

    # 优先用原始 bytes + BS4 自动编码检测，避免 httpx 误判编码
    raw_bytes = response.content
    soup = BeautifulSoup(raw_bytes, "html.parser")

    title = ""
    if soup.title and soup.title.string:
        title = soup.title.string.strip()
        # 清理标题中的网站后缀
        title = re.split(r"\s*[-_|–—]\s*", title)[0].strip()
    if not title:
        path = urlparse(url).path
        title = path.rstrip("/").rsplit("/", 1)[-1] or urlparse(url).netloc

    text = _extract_main_content(soup)

    if not text or len(text) < 50:
        raise RuntimeError(
            f"网页正文为空或过短（{len(text)} 字符），可能是 SPA 动态页面无法抓取: {url}"
        )

    return ParsedFile(
        title=title,
        content=text,
        file_type="html",
        source_path=url,
        metadata={"url": url, "status_code": response.status_code},
    )


def _parse_code(path: Path) -> ParsedFile:
    content = _read_file_text(path)
    lang = path.suffix.lstrip(".")
    header = f"文件: {path.name} (语言: {lang})\n{'=' * 40}\n"
    return ParsedFile(
        title=path.name,
        content=header + content,
        file_type="code",
        source_path=str(path),
        metadata={"language": lang},
    )


def _parse_image(path: Path) -> ParsedFile:
    from PIL import Image
    img = Image.open(str(path))
    metadata = {
        "format": img.format,
        "size": f"{img.width}x{img.height}",
        "mode": img.mode,
    }
    content = f"[图片文件: {path.name}]\n尺寸: {img.width}x{img.height}\n格式: {img.format or '未知'}"
    return ParsedFile(
        title=path.stem,
        content=content,
        file_type="image",
        source_path=str(path),
        metadata=metadata,
    )


def _parse_excel(path: Path) -> list[ParsedFile]:
    from openpyxl import load_workbook
    from openpyxl.utils.exceptions import InvalidFileException

    try:
        wb = load_workbook(str(path), data_only=True)
    except InvalidFileException:
        raise ValueError(
            f"不支持旧版 .xls 格式 ({path.name})，请先转换为 .xlsx 格式"
        )

    results = []
    sheet_names = list(wb.sheetnames)

    for name in sheet_names:
        ws = wb[name]

        merged_values = {}
        for merged_range in ws.merged_cells.ranges:
            top_left_val = ws.cell(row=merged_range.min_row, column=merged_range.min_col).value
            for r in range(merged_range.min_row, merged_range.max_row + 1):
                for c in range(merged_range.min_col, merged_range.max_col + 1):
                    if (r, c) != (merged_range.min_row, merged_range.min_col):
                        merged_values[(r, c)] = top_left_val

        raw_rows = []
        for row in ws.iter_rows(values_only=False):
            cells = []
            for cell in row:
                val = merged_values.get((cell.row, cell.column), cell.value)
                cells.append(str(val) if val is not None else "")
            raw_rows.append(cells)

        if not raw_rows:
            continue

        headers = _detect_headers(raw_rows[0])
        data_rows = raw_rows[1:] if headers else raw_rows

        # 结构化 Block 树（优先）
        structured = _table_rows_to_blocks(headers, data_rows, sheet_name=name)
        # 兼容：同时生成纯文本（用于全文搜索 fallback）
        text = _table_rows_to_text(headers, data_rows, sheet_name=name)
        if not text.strip():
            continue

        # 每个 sheet 独立知识条目，标题格式：文件名 - Sheet名
        title = f"{path.stem} - {name}" if len(sheet_names) > 1 else path.stem

        results.append(ParsedFile(
            title=title,
            content=text,
            file_type="xlsx",
            source_path=str(path),
            metadata={"sheet_name": name, "rows": len(data_rows), "sheets_total": len(sheet_names)},
            structured=structured,
        ))

    wb.close()

    # 无有效 sheet 时返回空列表（调用方应处理）
    if not results:
        results.append(ParsedFile(
            title=path.stem,
            content="",
            file_type="xlsx",
            source_path=str(path),
            metadata={"sheets": len(sheet_names)},
        ))

    return results


def _parse_csv(path: Path) -> ParsedFile:
    import csv
    import io

    text = _read_file_text(path)
    if not text.strip():
        return ParsedFile(
            title=path.stem, content="", file_type="csv",
            source_path=str(path), metadata={},
        )

    reader = csv.reader(io.StringIO(text))
    raw_rows = [[cell.strip() for cell in row] for row in reader]

    if not raw_rows:
        return ParsedFile(
            title=path.stem, content="", file_type="csv",
            source_path=str(path), metadata={},
        )

    headers = _detect_headers(raw_rows[0])
    data_rows = raw_rows[1:] if headers else raw_rows

    content = _table_rows_to_text(headers, data_rows, sheet_name="")
    structured = _table_rows_to_blocks(headers, data_rows, sheet_name="")

    return ParsedFile(
        title=path.stem,
        content=content,
        file_type="csv",
        source_path=str(path),
        metadata={"rows": len(data_rows)},
        structured=structured,
    )
